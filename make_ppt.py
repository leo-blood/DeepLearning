#!/usr/bin/env python3
"""Generate 分布式2026作业.pptx from motor_osdi24_zh.md"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── color palette ──────────────────────────────────────────
C_TITLE_BG  = RGBColor(0x1A, 0x37, 0x6C)   # deep navy
C_SLIDE_BG  = RGBColor(0xF4, 0xF7, 0xFF)   # very light blue
C_ACCENT    = RGBColor(0x2E, 0x75, 0xB6)   # azure
C_ACCENT2   = RGBColor(0xE8, 0x4C, 0x2E)   # red-orange
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK      = RGBColor(0x1A, 0x1A, 0x2E)
C_LIGHT_BOX = RGBColor(0xD6, 0xE4, 0xF7)
C_GRAY_TXT  = RGBColor(0x55, 0x55, 0x66)

W  = Inches(13.33)   # widescreen 16:9
H  = Inches(7.5)

# ── helpers ────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_layout(prs):
    return prs.slide_layouts[6]   # completely blank

def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def no_fill(shape):
    shape.fill.background()

def no_line(shape):
    shape.line.fill.background()

def add_rect(slide, x, y, w, h, color, radius=0):
    s = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    fill_solid(s, color)
    no_line(s)
    return s

def add_textbox(slide, text, x, y, w, h,
                font_size=20, bold=False, color=C_DARK,
                align=PP_ALIGN.LEFT, font_name="微软雅黑",
                wrap=True, word_wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return txb

def add_para(tf, text, font_size=18, bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, level=0, font_name="微软雅黑",
             space_before=6):
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return p

def slide_bg(slide, color=C_SLIDE_BG):
    """Paint background rectangle."""
    bg = add_rect(slide, 0, 0, W, H, color)
    # send to back
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

def header_bar(slide, title_text, sub_text="", bar_color=C_ACCENT, h=Inches(1.05)):
    bar = add_rect(slide, 0, 0, W, h, bar_color)
    add_textbox(slide, title_text,
                Inches(0.35), Inches(0.12), Inches(10), Inches(0.65),
                font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if sub_text:
        add_textbox(slide, sub_text,
                    Inches(0.35), Inches(0.62), Inches(10), Inches(0.38),
                    font_size=16, color=RGBColor(0xCC, 0xDD, 0xFF))

def footer_bar(slide, text="OSDI 2024 | Motor", color=C_ACCENT):
    bar = add_rect(slide, 0, H - Inches(0.32), W, Inches(0.32), color)
    add_textbox(slide, text,
                Inches(0.3), H - Inches(0.30), W - Inches(0.6), Inches(0.28),
                font_size=11, color=C_WHITE, align=PP_ALIGN.RIGHT)

def bullet_box(slide, items, x, y, w, h,
               bg_color=C_LIGHT_BOX, title=None, title_color=C_ACCENT,
               font_size=17, indent="▶  "):
    box = add_rect(slide, x, y, w, h, bg_color)
    no_line(box)
    tx_y = y + Inches(0.08)
    tx_h = h - Inches(0.1)
    if title:
        txb = slide.shapes.add_textbox(x + Inches(0.12), tx_y, w - Inches(0.2), Inches(0.4))
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(font_size + 2)
        run.font.bold = True
        run.font.name = "微软雅黑"
        run.font.color.rgb = title_color
        tx_y += Inches(0.38)
        tx_h -= Inches(0.38)
    txb = slide.shapes.add_textbox(x + Inches(0.15), tx_y, w - Inches(0.25), tx_h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = indent + item
        run.font.size = Pt(font_size)
        run.font.name = "微软雅黑"
        run.font.color.rgb = C_DARK
    return box

# ══════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════

def slide_title(prs):
    """Title slide"""
    s = prs.slides.add_slide(blank_layout(prs))
    # full-bleed background
    bg = add_rect(s, 0, 0, W, H, C_TITLE_BG)
    # decorative stripe
    stripe = add_rect(s, 0, H - Inches(1.6), W, Inches(1.6), RGBColor(0x0D, 0x1E, 0x45))
    # title
    add_textbox(s, "Motor：在分离式内存上为分布式事务启用多版本机制",
                Inches(0.7), Inches(1.6), Inches(11.9), Inches(2.2),
                font_size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # subtitle bar
    bar = add_rect(s, Inches(1.8), Inches(3.95), Inches(9.73), Inches(0.55),
                   RGBColor(0x2E, 0x75, 0xB6))
    add_textbox(s, "OSDI 2024  ·  Ming Zhang, Yu Hua, Zhijun Yang  ·  华中科技大学",
                Inches(1.9), Inches(3.98), Inches(9.5), Inches(0.5),
                font_size=17, color=C_WHITE, align=PP_ALIGN.CENTER)
    # meta
    add_textbox(s, "演讲人：[姓名]　　学号：[学号]",
                Inches(0.7), Inches(5.0), Inches(11.9), Inches(0.5),
                font_size=16, color=RGBColor(0xAA, 0xCC, 0xFF),
                align=PP_ALIGN.CENTER)
    add_textbox(s, "课程：分布式系统  ·  2026年6月",
                Inches(0.7), Inches(5.4), Inches(11.9), Inches(0.4),
                font_size=14, color=RGBColor(0x88, 0xAA, 0xDD),
                align=PP_ALIGN.CENTER)

def slide_outline(prs):
    s = prs.slides.add_slide(blank_layout(prs))
    slide_bg(s)
    header_bar(s, "目录", "Outline")
    footer_bar(s)
    items = [
        ("01", "研究背景与动机"),
        ("02", "挑战：多版本机制的移植困境"),
        ("03", "Motor 总体架构"),
        ("04", "连续版本元组（CVT）"),
        ("05", "独立值域与内存优化"),
        ("06", "垃圾回收与锚标志一致性"),
        ("07", "MVCC 事务协议"),
        ("08", "实验评估"),
        ("09", "结论与总结"),
    ]
    cols = 3
    rows = 3
    bw = Inches(3.8)
    bh = Inches(1.5)
    gx = Inches(0.35)
    gy = Inches(1.18)
    gap_x = Inches(0.2)
    gap_y = Inches(0.18)
    for idx, (num, title) in enumerate(items):
        col = idx % cols
        row = idx // cols
        x = gx + col * (bw + gap_x)
        y = gy + row * (bh + gap_y)
        box = add_rect(s, x, y, bw, bh, C_ACCENT)
        no_line(box)
        add_textbox(s, num, x + Inches(0.15), y + Inches(0.12), Inches(0.7), Inches(0.55),
                    font_size=30, bold=True, color=RGBColor(0xFF, 0xDD, 0x44))
        add_textbox(s, title, x + Inches(0.15), y + Inches(0.72), bw - Inches(0.25), Inches(0.65),
                    font_size=17, bold=True, color=C_WHITE)

def slide_background(prs):
    """内存分离背景"""
    s = prs.slides.add_slide(blank_layout(prs))
    slide_bg(s)
    header_bar(s, "研究背景：内存分离架构", "Memory Disaggregation")
    footer_bar(s)

    # Left box: 传统 vs 分离
    bullet_box(s,
        ["传统单体服务器：计算+内存紧耦合，资源利用率低",
         "故障域粗糙：一 CPU 故障导致整机不可用",
         "内存分离：将计算池与内存池解耦，通过 RDMA/CXL 互连",
         "计算池：强 CPU，少量 DRAM 用于元数据缓存",
         "内存池：大容量 DRAM，弱计算单元（仅内存分配）"],
        Inches(0.3), Inches(1.18), Inches(6.1), Inches(3.7),
        title="传统单体 vs 内存分离", font_size=16)

    # Right box: 优势
    bullet_box(s,
        ["资源利用率显著提升（按需分配）",
         "弹性扩缩容（独立扩展计算/内存）",
         "故障隔离（故障域缩小至单资源池）",
         "代表性互连技术：RDMA（100 Gbps）、CXL"],
        Inches(6.6), Inches(1.18), Inches(6.4), Inches(2.4),
        bg_color=RGBColor(0xE0, 0xF0, 0xE8),
        title="内存分离的优势", font_size=16,
        title_color=RGBColor(0x1A, 0x7A, 0x3A))

    # RDMA box
    bullet_box(s,
        ["单侧 RDMA：绕过远程 CPU 直接读写内存",
         "支持 READ / WRITE / CAS / FAA 原子操作",
         "100 Gbps InfiniBand，微秒级延迟"],
        Inches(6.6), Inches(3.75), Inches(6.4), Inches(2.0),
        bg_color=RGBColor(0xFF, 0xF3, 0xE0),
        title="单侧 RDMA 操作", font_size=16,
        title_color=RGBColor(0xB8, 0x5A, 0x00))

    add_textbox(s, "核心挑战：如何在弱计算内存池上实现高效分布式事务？",
                Inches(0.3), Inches(5.05), Inches(12.7), Inches(0.38),
                font_size=17, bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)

def slide_motivation(prs):
    """单版本设计局限"""
    s = prs.slides.add_slide(blank_layout(prs))
    slide_bg(s)
    header_bar(s, "研究动机：单版本设计的局限性", "Motivation")
    footer_bar(s)

    bullet_box(s,
        ["FORD：分离式内存上单版本事务代表系统",
         "每条数据仅维护一个版本（最新版本）"],
        Inches(0.3), Inches(1.18), Inches(12.7), Inches(1.1),
        bg_color=RGBColor(0xE8, 0xEE, 0xF8), title=None, font_size=17)

    # two pain points
    add_rect(s, Inches(0.3), Inches(2.45), Inches(6.1), Inches(2.5),
             RGBColor(0xFD, 0xE8, 0xE8))
    add_textbox(s, "问题 ①  并发度低",
                Inches(0.45), Inches(2.52), Inches(5.8), Inches(0.45),
                font_size=18, bold=True, color=C_ACCENT2)
    add_textbox(s,
        "事务提交期间，正在更新的数据不可读，\n"
        "读操作必须等待写操作完成，形成阻塞。\n"
        "→ 高并发下吞吐量严重受限",
        Inches(0.45), Inches(3.02), Inches(5.8), Inches(1.8),
        font_size=16, color=C_DARK)

    add_rect(s, Inches(6.7), Inches(2.45), Inches(6.3), Inches(2.5),
             RGBColor(0xFD, 0xE8, 0xE8))
    add_textbox(s, "问题 ②  Undo 日志开销高",
                Inches(6.85), Inches(2.52), Inches(6.0), Inches(0.45),
                font_size=18, bold=True, color=C_ACCENT2)
    add_textbox(s,
        "为保证原子性，需向所有副本写入 undo 日志，\n"
        "消耗网络带宽，协调者需等待所有 ACK，\n"
        "→ 提交延迟高、吞吐量低",
        Inches(6.85), Inches(3.02), Inches(6.0), Inches(1.8),
        font_size=16, color=C_DARK)

    add_textbox(s, "解决方案：引入多版本机制（MVCC）",
                Inches(0.3), Inches(5.12), Inches(12.7), Inches(0.38),
                font_size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s,
        "✓ 读请求获取现有版本，无需等待写操作  →  提升并发   "
        "✓ 旧版本天然充当 undo 日志  →  消除日志写入开销",
        Inches(0.3), Inches(5.55), Inches(12.7), Inches(0.5),
        font_size=16, color=C_DARK, align=PP_ALIGN.CENTER)

def slide_challenges(prs):
    """移植挑战"""
    s = prs.slides.add_slide(blank_layout(prs))
    slide_bg(s)
    header_bar(s, "挑战：多版本机制无法直接移植到分离式内存", "Challenges")
    footer_bar(s)

    add_textbox(s, "现有多版本系统均面向传统单体架构设计，移植至分离式内存面临两大核心挑战：",
                Inches(0.35), Inches(1.22), Inches(12.6), Inches(0.4),
                font_size=17, color=C_DARK)

    # challenge 1
    add_rect(s, Inches(0.3), Inches(1.75), Inches(12.7), Inches(2.1), C_LIGHT_BOX)
    add_textbox(s, "挑战 1：事务协议不兼容",
                Inches(0.45), Inches(1.82), Inches(12.0), Inches(0.42),
                font_size=19, bold=True, color=C_ACCENT)
    add_textbox(s,
        "现有系统假设每台服务器拥有强大 CPU 执行协议中大量计算任务（加锁、验证、时间戳计算等）。\n"
        "但内存池中 CPU 计算能力极弱，无法频繁处理这些任务。\n"
        "→ 需要设计全新的「单侧 RDMA」事务协议，绕过内存池弱 CPU",
        Inches(0.45), Inches(2.28), Inches(12.3), Inches(1.4),
        font_size=16, color=C_DARK)

    # challenge 2
    add_rect(s, Inches(0.3), Inches(3.98), Inches(12.7), Inches(2.55), C_LIGHT_BOX)
    add_textbox(s, "挑战 2：链式版本结构效率低下",
                Inches(0.45), Inches(4.05), Inches(12.0), Inches(0.42),
                font_size=19, bold=True, color=C_ACCENT)
    add_textbox(s,
        "现有方案（旧→新链 / 新→旧链）：读取特定版本需沿指针逐个拉取 → 每步都是一次网络往返（RTT）\n"
        "实测：链式遍历步数从 1 增加到 20，RDMA 读取延迟增加 24.8 倍\n"
        "GC（垃圾回收）需追踪最旧运行事务，但内存池计算单元无法感知事务状态 → GC 开销极大",
        Inches(0.45), Inches(4.52), Inches(12.3), Inches(1.85),
        font_size=16, color=C_DARK)

def slide_overview(prs):
    """Motor 总体架构"""
    s = prs.slides.add_slide(blank_layout(prs))
    slide_bg(s)
    header_bar(s, "Motor 总体架构", "System Overview")
    footer_bar(s)

    bullet_box(s,
        ["Motor 内存存储（Memory Storage）：CVT 结构，高效组织多版本数据",
         "Motor 事务协议（Transaction Protocol）：完全单侧 RDMA，支持多隔离级别",
         "两者协同：协调者通过 RDMA 绕过内存池弱 CPU，实现全流程处理"],
        Inches(0.3), Inches(1.18), Inches(12.7), Inches(1.45),
        font_size=16)

    # workflow steps
    steps = [
        ("① 数据加载", "客户端通过内存池 CPU\n分配内存，加载 DB 表\n（CVT 结构 + 索引）"),
        ("② 建立连接", "计算池与内存池建立\nRDMA 连接，内存池\n发送元数据给协调者"),
        ("③ 事务请求", "客户端向计算池\n发出事务请求"),
        ("④ 协调者执行", "协调者并发运行 Motor\n协议：锁定、验证、\n提交（纯单侧 RDMA）"),
    ]
    sw = Inches(2.9)
    sh = Inches(2.4)
    sy = Inches(2.82)
    for i, (title, body) in enumerate(steps):
        sx = Inches(0.35) + i * (sw + Inches(0.3))
        box = add_rect(s, sx, sy, sw, sh, C_ACCENT)
        add_textbox(s, title, sx + Inches(0.1), sy + Inches(0.12),
                    sw - Inches(0.15), Inches(0.45),
                    font_size=17, bold=True, color=RGBColor(0xFF, 0xEE, 0x88))
        add_textbox(s, body, sx + Inches(0.1), sy + Inches(0.62),
                    sw - Inches(0.15), Inches(1.6),
                    font_size=15, color=C_WHITE)
        if i < 3:
            add_textbox(s, "→", sx + sw + Inches(0.04), sy + Inches(0.95),
                        Inches(0.3), Inches(0.5),
                        font_size=24, bold=True, color=C_ACCENT)

    add_textbox(s, "关键设计原则：内存池不执行计算，所有计算逻辑均在计算池协调者中完成",
                Inches(0.3), Inches(5.38), Inches(12.7), Inches(0.38),
                font_size=16, bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)

def slide_cvt_design(prs):
    """CVT 设计"""
    s = prs.slides.add_slide(blank_layout(prs))
    slide_bg(s)
    header_bar(s, "连续版本元组（CVT）结构", "Consecutive Version Tuple")
    footer_bar(s)

    add_textbox(s, "核心思想：放弃指针链接，将多个版本连续存储于连续地址空间",
                Inches(0.35), Inches(1.22), Inches(12.6), Inches(0.38),
                font_size=17, bold=True, color=C_ACCENT)

    # Left: structure description
    bullet_box(s,
        ["Header：TableID / Key / Lock / AttrBarPtr / VpkgPtr",
         "Vcell × VNum：VcellSA + VcellEA（锚）/ Valid / Version / Bitmap / StartOffset",
         "一次 RDMA READ 获取整个 CVT（含所有版本）",
         "本地搜索目标版本，无网络 I/O → 极低延迟"],
        Inches(0.3), Inches(1.72), Inches(6.1), Inches(3.5),
        title="CVT 结构字段", font_size=15)

    # Right: comparison
    add_rect(s, Inches(6.6), Inches(1.72), Inches(6.4), Inches(3.5),
             RGBColor(0xFF, 0xF8, 0xE8))
    add_textbox(s, "链式结构 vs CVT 对比",
                Inches(6.75), Inches(1.80), Inches(6.1), Inches(0.42),
                font_size=17, bold=True, color=RGBColor(0xB8, 0x6A, 0x00))

    rows_data = [
        ("对比项", "旧→新 / 新→旧链", "CVT（Motor）"),
        ("读取往返次数", "多次（链式遍历）", "1 次"),
        ("延迟（增至20步）", "增加 24.8 倍", "不变"),
        ("GC 方式", "需追踪事务状态", "协调者主动抢占"),
        ("内存结构", "动态指针链接", "连续地址空间"),
    ]
    col_w = [Inches(1.8), Inches(2.1), Inches(2.0)]
    row_h = Inches(0.46)
    tx = Inches(6.65)
    ty = Inches(2.28)
    for ri, row in enumerate(rows_data):
        for ci, cell in enumerate(row):
            cx = tx + sum(col_w[:ci])
            bg = C_ACCENT if ri == 0 else (C_LIGHT_BOX if ri % 2 == 0 else C_WHITE)
            fc = C_WHITE if ri == 0 else C_DARK
            add_rect(s, cx, ty + ri * row_h, col_w[ci], row_h, bg)
            add_textbox(s, cell, cx + Inches(0.05), ty + ri * row_h + Inches(0.06),
                        col_w[ci] - Inches(0.08), row_h - Inches(0.08),
                        font_size=13, bold=(ri == 0), color=fc, align=PP_ALIGN.CENTER)

    add_textbox(s, "VNum 配置：低争用短事务（TATP）→ VNum=2；高争用长事务（TPCC）→ VNum=4",
                Inches(0.3), Inches(5.42), Inches(12.7), Inches(0.38),
                font_size=15, bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)
    add_textbox(s, "支持哈希表 / B+ 树索引统一接口；CVT 地址缓存避免重复拉取哈希桶",
                Inches(0.3), Inches(5.85), Inches(12.7), Inches(0.35),
                font_size=14, color=C_GRAY_TXT, align=PP_ALIGN.CENTER)

def slide_value_region(prs):
    """独立值域"""
    s = prs.slides.add_slide(blank_layout(prs))
    slide_bg(s)
    header_bar(s, "独立值域与内存开销优化", "Independent Value Region")
    footer_bar(s)

    bullet_box(s,
        ["CVT（版本元数据）与数据值分开存储",
         "先读 CVT 确定目标版本 → 再读对应值（按需传输）",
         "CVT 大小稳定，不受值大小影响 → 读取延迟可预测"],
        Inches(0.3), Inches(1.18), Inches(12.7), Inches(1.2),
        font_size=16)

    # two areas
    add_rect(s, Inches(0.3), Inches(2.55), Inches(6.1), Inches(2.2), RGBColor(0xE8, 0xF4, 0xE8))
    add_textbox(s, "全值区（Full-value Area）",
                Inches(0.45), Inches(2.62), Inches(5.8), Inches(0.42),
                font_size=17, bold=True, color=RGBColor(0x1A, 0x7A, 0x3A))
    add_textbox(s,
        "存储最新版本的完整值\n"
        "直接用一次 RDMA READ 读取（最新版本）",
        Inches(0.45), Inches(3.08), Inches(5.8), Inches(1.5),
        font_size=15, color=C_DARK)

    add_rect(s, Inches(6.6), Inches(2.55), Inches(6.4), Inches(2.2), RGBColor(0xFF, 0xF3, 0xE0))
    add_textbox(s, "Delta 区（属性条带）",
                Inches(6.75), Inches(2.62), Inches(6.1), Inches(0.42),
                font_size=17, bold=True, color=RGBColor(0xB8, 0x5A, 0x00))
    add_textbox(s,
        "仅存储被修改的旧属性（非完整值）\n"
        "Bitmap + StartOffset 定位各版本属性偏移\n"
        "构造旧版本：完整值覆盖旧属性即可",
        Inches(6.75), Inches(3.08), Inches(6.1), Inches(1.5),
        font_size=15, color=C_DARK)

    bullet_box(s,
        ["一次 RTT 读取：批量 RDMA READ 同时读完整值 + 旧属性，本地重构目标版本",
         "一次 RTT 写入：批量 RDMA WRITE 同时更新完整值 + 追加旧属性到属性条带",
         "内存节省：存修改属性而非完整值；TPCC 支持 4 版本，内存仅增至 FORD 的 1.45 倍（非 4 倍）"],
        Inches(0.3), Inches(4.92), Inches(12.7), Inches(1.65),
        title="关键优化", font_size=15)

def slide_gc_anchor(prs):
    """GC 与锚标志"""
    s = prs.slides.add_slide(blank_layout(prs))
    slide_bg(s)
    header_bar(s, "垃圾回收 & 锚标志辅助一致性", "GC & Anchor-assisted Reading")
    footer_bar(s)

    # GC
    add_rect(s, Inches(0.3), Inches(1.18), Inches(6.1), Inches(3.85), C_LIGHT_BOX)
    add_textbox(s, "协调者主动 GC",
                Inches(0.45), Inches(1.25), Inches(5.8), Inches(0.42),
                font_size=19, bold=True, color=C_ACCENT)
    add_textbox(s,
        "传统 GC：需追踪最旧运行事务 → 内存池无能力\n\n"
        "Motor 方案：\n"
        "• CVT 无空闲 Vcell 时，协调者主动抢占\n"
        "• 选择「最旧版本」作为受害版本直接覆盖\n"
        "• 无需追踪任何事务状态\n\n"
        "RDMA 加速事务 → 最旧版本仍在使用概率极低\n"
        "合适 VNum 可高效缓解少数长事务中止",
        Inches(0.45), Inches(1.72), Inches(5.8), Inches(3.15),
        font_size=15, color=C_DARK)

    # Anchor
    add_rect(s, Inches(6.6), Inches(1.18), Inches(6.4), Inches(3.85), RGBColor(0xFF, 0xF8, 0xE8))
    add_textbox(s, "锚标志辅助读取（4个锚：各 1 字节）",
                Inches(6.75), Inches(1.25), Inches(6.1), Inches(0.42),
                font_size=17, bold=True, color=RGBColor(0xB8, 0x6A, 0x00))
    add_textbox(s,
        "问题：读取 CVT 时，另一协调者可能正在执行 GC\n"
        "→ 读到被部分更新的损坏值\n\n"
        "四个锚：VcellSA / VcellEA / VpkgSA / VpkgEA\n\n"
        "写入顺序：\n"
        "  先写 Vpkg → 写修改属性 → 写 Vcell\n"
        "  四个锚同时 +1 保持相等\n\n"
        "读取验证：\n"
        "  若 4 锚相等 → 值未被修改 → 安全重构\n"
        "  若不相等   → 检测到冲突 → 事务中止",
        Inches(6.75), Inches(1.72), Inches(6.1), Inches(3.15),
        font_size=14, color=C_DARK)

    add_textbox(s, "保证写入顺序：禁用内存池 DDIO，确保 RDMA WRITE 以 FIFO 顺序写入主存",
                Inches(0.3), Inches(5.2), Inches(12.7), Inches(0.35),
                font_size=14, bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)

def slide_protocol(prs):
    """事务协议"""
    s = prs.slides.add_slide(blank_layout(prs))
    slide_bg(s)
    header_bar(s, "MVCC 事务协议：三阶段处理", "Transaction Protocol")
    footer_bar(s)

    phases = [
        ("阶段 1\n执行（Execution）",
         "• 获取 T_start（开始时间戳）\n"
         "• 查 CVT 地址缓存\n"
         "• RW 数据：CAS+READ 一次 RTT 加锁并读 CVT\n"
         "• RO 数据：READ 读 CVT\n"
         "• 选目标版本 V0（≤T_start 最大版本）\n"
         "• 提前中止：若存在 >T_start 的版本\n"
         "• 批量 READ 取值，检查 4 个锚标志",
         C_ACCENT, "1 RTT（加锁+读）"),
        ("阶段 2\n验证（Validation）",
         "• 获取 T_commit（提交时间戳）\n"
         "• 仅含 RW 数据 → 跳过验证\n"
         "• 重新读取 RO 数据 CVT\n"
         "• 检查是否被锁或出现新版本 V'≠V0\n"
         "• 任一条件成立 → 中止事务",
         RGBColor(0x1A, 0x6A, 0x3A), "1 RTT（验证 RO）"),
        ("阶段 3\n提交（Commit）",
         "• 准备新 Vcell + 新 Vpkg（本地完成）\n"
         "• GC（若 CVT/属性条带已满）\n"
         "• 批量 RDMA WRITE 一次 RTT\n"
         "  写入所有副本 + 解锁主副本\n"
         "• 收到所有 ACK → 向应用报告已提交",
         RGBColor(0x8B, 0x00, 0x00), "1 RTT（提交+解锁）"),
    ]

    pw = Inches(3.85)
    ph = Inches(4.18)
    py = Inches(1.22)
    for i, (title, body, col, rtt) in enumerate(phases):
        px = Inches(0.25) + i * (pw + Inches(0.27))
        box = add_rect(s, px, py, pw, ph, col)
        add_textbox(s, title, px + Inches(0.12), py + Inches(0.1),
                    pw - Inches(0.2), Inches(0.75),
                    font_size=17, bold=True, color=C_WHITE)
        add_textbox(s, body, px + Inches(0.12), py + Inches(0.9),
                    pw - Inches(0.2), ph - Inches(1.1),
                    font_size=14, color=C_WHITE)
        # RTT badge
        badge = add_rect(s, px + pw - Inches(1.7), py + ph - Inches(0.42),
                         Inches(1.65), Inches(0.38), RGBColor(0xFF, 0xDD, 0x44))
        add_textbox(s, rtt, px + pw - Inches(1.68), py + ph - Inches(0.40),
                    Inches(1.6), Inches(0.35),
                    font_size=12, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        if i < 2:
            add_textbox(s, "→", px + pw + Inches(0.04), py + ph / 2 - Inches(0.3),
                        Inches(0.25), Inches(0.5), font_size=22, bold=True, color=C_ACCENT)

    add_textbox(s, "只读事务：无需验证，直接获取 T_start 时刻快照 → 最低延迟",
                Inches(0.25), Inches(5.55), Inches(12.8), Inches(0.35),
                font_size=15, bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)

def slide_isolation(prs):
    """隔离级别"""
    s = prs.slides.add_slide(blank_layout(prs))
    slide_bg(s)
    header_bar(s, "灵活支持隔离级别 & 容错机制", "Isolation Levels & Fault Tolerance")
    footer_bar(s)

    # Isolation
    add_rect(s, Inches(0.3), Inches(1.18), Inches(6.1), Inches(3.85), C_LIGHT_BOX)
    add_textbox(s, "两种隔离级别",
                Inches(0.45), Inches(1.25), Inches(5.8), Inches(0.42),
                font_size=19, bold=True, color=C_ACCENT)
    add_textbox(s,
        "可串行化（SR）\n"
        "  加锁保证 RW 数据版本不变\n"
        "  验证保证 RO 数据版本不变\n"
        "  只读事务移至合适串行化执行点\n\n"
        "快照隔离（SI）\n"
        "  禁用 RO 数据版本验证\n"
        "  允许读取 T_start 时刻旧快照\n"
        "  仍加锁解决写写冲突\n"
        "  MySQL / PostgreSQL / Oracle 均采用\n\n"
        "完整 ACID 保证：多版本→原子性，\n"
        "副本复制→持久性，锁+验证→隔离性",
        Inches(0.45), Inches(1.72), Inches(5.8), Inches(3.15),
        font_size=14, color=C_DARK)

    # Fault tolerance
    add_rect(s, Inches(6.6), Inches(1.18), Inches(6.4), Inches(3.85), RGBColor(0xF0, 0xF4, 0xFF))
    add_textbox(s, "容错机制",
                Inches(6.75), Inches(1.25), Inches(6.1), Inches(0.42),
                font_size=19, bold=True, color=C_ACCENT)
    add_textbox(s,
        "内存池副本故障\n"
        "  提交前故障 → 协调者丢弃并中止\n"
        "  主副本故障 → 备副本提升为主\n"
        "  备副本故障 → 新节点替换，数据迁移\n\n"
        "协调者故障（租约检测）\n"
        "  操作日志最大 556B/事务（本地内存）\n"
        "  新协调者利用日志恢复飞行中提交\n"
        "  约 170ms 恢复至峰值吞吐量\n\n"
        "网络故障\n"
        "  主分区继续服务（牺牲部分可用性）\n"
        "  参照 uKharon 由管理员解决分区",
        Inches(6.75), Inches(1.72), Inches(6.1), Inches(3.15),
        font_size=14, color=C_DARK)

def slide_eval_setup(prs):
    """实验配置"""
    s = prs.slides.add_slide(blank_layout(prs))
    slide_bg(s)
    header_bar(s, "实验评估：配置与基准", "Experimental Setup")
    footer_bar(s)

    bullet_box(s,
        ["4 台服务器，100Gbps InfiniBand（Mellanox SB7890）互连",
         "计算池：1 台，Intel Xeon Gold 6330 CPU",
         "内存池：3 台，每台 192GB DRAM",
         "3 路副本（1 主 2 备）"],
        Inches(0.3), Inches(1.18), Inches(12.7), Inches(1.45),
        title="实验平台", font_size=15)

    benchmarks = [
        ("KVS\n微基准", "10M 键值对\n键 8B / 值 40B\nZipfian 分布"),
        ("TATP\n电信应用", "4 张 DB 表\n80% 只读事务\n2M 用户 / 最大 48B"),
        ("SmallBank\n银行应用", "2 张 DB 表\n85% 读写事务\n10M 账户 / 16B"),
        ("TPCC\n复杂订单", "9 张 DB 表\n92% 读写事务\n24 仓库 / 最大 672B"),
    ]
    bw = Inches(2.9)
    bh = Inches(1.6)
    by = Inches(2.82)
    colors = [C_ACCENT, RGBColor(0x1A,0x7A,0x3A), RGBColor(0xB8,0x5A,0x00), C_ACCENT2]
    for i, (name, detail) in enumerate(benchmarks):
        bx = Inches(0.35) + i * (bw + Inches(0.3))
        add_rect(s, bx, by, bw, bh, colors[i])
        add_textbox(s, name, bx + Inches(0.12), by + Inches(0.1),
                    bw - Inches(0.2), Inches(0.62),
                    font_size=16, bold=True, color=C_WHITE)
        add_textbox(s, detail, bx + Inches(0.12), by + Inches(0.75),
                    bw - Inches(0.2), Inches(0.75),
                    font_size=13, color=C_WHITE)

    bullet_box(s,
        ["FaRMv2-DM：FaRMv2（单体多版本）用单侧 RDMA 重新实现 → 兼容分离式内存",
         "FORD：分离式内存单版本事务代表系统"],
        Inches(0.3), Inches(4.58), Inches(12.7), Inches(1.0),
        title="对比系统", font_size=15, title_color=C_ACCENT2,
        bg_color=RGBColor(0xFD, 0xE8, 0xE8))

def slide_eval_main(prs):
    """端到端性能"""
    s = prs.slides.add_slide(blank_layout(prs))
    slide_bg(s)
    header_bar(s, "实验结果：端到端性能", "End-to-End Performance")
    footer_bar(s)

    add_textbox(s, "Motor vs FORD（单版本）",
                Inches(0.3), Inches(1.22), Inches(6.1), Inches(0.42),
                font_size=18, bold=True, color=C_ACCENT)
    rows1 = [
        ("基准", "吞吐量提升", "P50 延迟降低"),
        ("TATP", "+14.4%", "FORD 更低*"),
        ("TPCC", "+98.1%", "–55.8%"),
        ("SmallBank", "+65.4%", "–26.2%"),
    ]
    cw1 = [Inches(1.5), Inches(2.0), Inches(2.1)]
    rh = Inches(0.44)
    tx, ty = Inches(0.3), Inches(1.72)
    for ri, row in enumerate(rows1):
        for ci, cell in enumerate(row):
            cx = tx + sum(cw1[:ci])
            bg = C_ACCENT if ri == 0 else (C_LIGHT_BOX if ri % 2 == 0 else C_WHITE)
            fc = C_WHITE if ri == 0 else C_DARK
            add_rect(s, cx, ty + ri * rh, cw1[ci], rh, bg)
            add_textbox(s, cell, cx + Inches(0.04), ty + ri * rh + Inches(0.06),
                        cw1[ci] - Inches(0.06), rh - Inches(0.08),
                        font_size=13, bold=(ri == 0), color=fc, align=PP_ALIGN.CENTER)

    add_textbox(s, "*TATP 中 70% 事务只读 1 对象，FORD 仅 1 RTT；Motor 需 2 RTT（CVT+值）",
                Inches(0.3), Inches(3.6), Inches(6.1), Inches(0.38),
                font_size=12, color=C_GRAY_TXT)

    add_textbox(s, "Motor vs FaRMv2-DM（多版本）",
                Inches(6.6), Inches(1.22), Inches(6.4), Inches(0.42),
                font_size=18, bold=True, color=C_ACCENT)
    rows2 = [
        ("基准", "吞吐量提升", "P50/P99 延迟降低"),
        ("TATP", "+18.9%", "8.6% / 39.1%"),
        ("TPCC", "+44.3%", "52.1% / 35.6%"),
        ("SmallBank", "+29.5%", "43.6% / 34.5%"),
    ]
    cw2 = [Inches(1.3), Inches(1.8), Inches(2.7)]
    tx2 = Inches(6.6)
    for ri, row in enumerate(rows2):
        for ci, cell in enumerate(row):
            cx = tx2 + sum(cw2[:ci])
            bg = C_ACCENT if ri == 0 else (C_LIGHT_BOX if ri % 2 == 0 else C_WHITE)
            fc = C_WHITE if ri == 0 else C_DARK
            add_rect(s, cx, ty + ri * rh, cw2[ci], rh, bg)
            add_textbox(s, cell, cx + Inches(0.04), ty + ri * rh + Inches(0.06),
                        cw2[ci] - Inches(0.06), rh - Inches(0.08),
                        font_size=13, bold=(ri == 0), color=fc, align=PP_ALIGN.CENTER)

    bullet_box(s,
        ["CVT 一次往返获取所有版本，FaRMv2 链式结构需多次往返",
         "Motor 批量发送加锁+读取请求，FaRMv2 需专用 RTT 加锁",
         "Motor 一次 RTT 提交所有副本，FaRMv2 需两次 RTT"],
        Inches(0.3), Inches(3.85), Inches(12.7), Inches(1.5),
        title="Motor 优于 FaRMv2-DM 的三大原因", font_size=14,
        title_color=RGBColor(0x1A, 0x7A, 0x3A),
        bg_color=RGBColor(0xE8, 0xF4, 0xE8))

def slide_eval_memory(prs):
    """内存 & 隔离"""
    s = prs.slides.add_slide(blank_layout(prs))
    slide_bg(s)
    header_bar(s, "内存开销 & 隔离级别性能", "Memory Usage & Isolation")
    footer_bar(s)

    bullet_box(s,
        ["FORD（单版本）：内存最低（基准线）",
         "Motor（4版本 TPCC）：内存仅增至 FORD 的 1.45 倍（非 4 倍）",
         "FaRMv2-DM：比 Motor 高 14.6%–22.8%（每版本存完整值 + 指针开销）",
         "Motor 节省内存三策略：① 存修改属性 ② 精确估算 ABS ③ 按负载配置 VNum",
         "VNum 从 2→8（+4 倍），内存占用仅增约 1.4–2.1 倍（远低于线性）"],
        Inches(0.3), Inches(1.18), Inches(12.7), Inches(2.15),
        title="内存开销分析", font_size=15)

    # SI vs SR
    add_rect(s, Inches(0.3), Inches(3.48), Inches(12.7), Inches(1.88),
             RGBColor(0xFF, 0xF3, 0xE0))
    add_textbox(s, "隔离级别对性能的影响（Motor-SR vs Motor-SI）",
                Inches(0.45), Inches(3.55), Inches(12.3), Inches(0.42),
                font_size=18, bold=True, color=RGBColor(0xB8, 0x5A, 0x00))
    add_textbox(s,
        "Motor-SI（快照隔离）消除 RO 数据验证阶段，在读密集（TATP）和写密集（TPCC）上均优于 Motor-SR\n"
        "TPCC 改善更显著：每事务访问更多 RO 数据 + 更高读写争用，放宽隔离带来更大收益\n"
        "已被 MySQL / PostgreSQL / Oracle / SQL Server 采用，在实际 OLTP 中广泛可用",
        Inches(0.45), Inches(4.02), Inches(12.3), Inches(1.25),
        font_size=14, color=C_DARK)

    bullet_box(s,
        ["内存池使用 Intel Optane PM（持久内存），吞吐量仅降低 13.1%",
         "Motor 在 DRAM 和 PM 上均高效工作，具备良好可移植性"],
        Inches(0.3), Inches(5.5), Inches(12.7), Inches(0.88),
        title="PM 扩展性", font_size=14,
        bg_color=RGBColor(0xE8, 0xEE, 0xF8), title_color=C_ACCENT)

def slide_conclusion(prs):
    s = prs.slides.add_slide(blank_layout(prs))
    slide_bg(s)
    header_bar(s, "结论与总结", "Conclusion")
    footer_bar(s)

    contribs = [
        ("CVT 结构",
         "连续存储多版本，一次 RDMA READ 获取所有版本，消除链式遍历多 RTT 开销"),
        ("协调者主动 GC",
         "无需追踪事务状态，抢占式覆盖最旧版本，轻量高效"),
        ("锚标志一致性",
         "4 个锚标志检测并发 GC 导致的部分更新，保证值读取正确性"),
        ("单侧 RDMA 协议",
         "三阶段（执行/验证/提交）完全绕过内存池弱 CPU，支持 SR 和 SI 隔离级别"),
        ("显著性能提升",
         "吞吐量最高 +98.1%，延迟最低 –55.8%，内存开销适度（TPCC 仅 1.45× FORD）"),
    ]

    bw = Inches(4.9)
    bh = Inches(0.88)
    gap = Inches(0.14)
    for i, (title, body) in enumerate(contribs):
        col = i % 2
        row = i // 2
        if i == 4:
            bx = (W - bw) / 2
        else:
            bx = Inches(0.3) + col * (bw + Inches(0.8))
        by = Inches(1.22) + row * (bh + gap)
        color = [C_ACCENT, RGBColor(0x1A,0x7A,0x3A),
                 RGBColor(0xB8,0x5A,0x00), C_ACCENT2,
                 RGBColor(0x55, 0x33, 0x99)][i]
        add_rect(s, bx, by, bw, bh, color)
        add_textbox(s, title, bx + Inches(0.12), by + Inches(0.06),
                    bw - Inches(0.2), Inches(0.38),
                    font_size=16, bold=True, color=C_WHITE)
        add_textbox(s, body, bx + Inches(0.12), by + Inches(0.46),
                    bw - Inches(0.2), Inches(0.38),
                    font_size=13, color=C_WHITE)

    add_textbox(s,
        "Motor 开源地址：https://github.com/minghust/motor  ·  "
        "发表于 OSDI 2024（第 801–819 页）",
        Inches(0.3), Inches(6.85), Inches(12.7), Inches(0.35),
        font_size=13, color=C_GRAY_TXT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════

def build_ppt(out_path):
    prs = new_prs()
    slide_title(prs)
    slide_outline(prs)
    slide_background(prs)
    slide_motivation(prs)
    slide_challenges(prs)
    slide_overview(prs)
    slide_cvt_design(prs)
    slide_value_region(prs)
    slide_gc_anchor(prs)
    slide_protocol(prs)
    slide_isolation(prs)
    slide_eval_setup(prs)
    slide_eval_main(prs)
    slide_eval_memory(prs)
    slide_conclusion(prs)
    prs.save(out_path)
    print(f"Saved {len(prs.slides)} slides → {out_path}")

if __name__ == "__main__":
    out = "/Users/wen/Desktop/未命名文件夹/DeepLearning/学号_姓名_分布式2026作业.pptx"
    build_ppt(out)
